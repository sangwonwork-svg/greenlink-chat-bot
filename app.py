import streamlit as st
from groq import Groq
import fitz  # PDF용
from pptx import Presentation  # PPTX용
import pandas as pd  # Excel용
import os
import olefile # HWP용

# --- 1. 보안 설정 ---
GROQ_API_KEY = st.secrets["GROQ_API_KEY"]
ACCESS_PASSWORD = st.secrets["ACCESS_PASSWORD"]

st.set_page_config(page_title="사내 통합 지식고", layout="wide")

# --- 2. 로그인 기능 ---
if "auth" not in st.session_state:
    st.session_state.auth = False

if not st.session_state.auth:
    st.title("🔒 지식고 로그인")
    pwd = st.text_input("비밀번호 입력:", type="password")
    if st.button("접속"):
        if pwd == ACCESS_PASSWORD:
            st.session_state.auth = True
            st.rerun()
        else:
            st.error("비밀번호가 틀렸습니다.")
    st.stop()

# --- 3. 파일 통합 로직 ---
@st.cache_resource
def load_all_documents():
    combined_text = ""
    file_list = []
    
    # 순수하게 '학습용 문서'로 간주할 확장자만 지정
    target_extensions = (".pdf", ".pptx", ".txt", ".hwp", ".xlsx", ".xls")
    # 학습에서 제외할 특정 파일명 지정
    exclude_files = ("requirements.txt", "app.py", "packages.txt")
    
    for file in os.listdir("."):
        # 확장자가 대상에 포함되고, 제외 목록에 없으며, 숨김 파일(.)이 아닌 경우만 처리
        if file.lower().endswith(target_extensions) and file not in exclude_files and not file.startswith("."):
            try:
                # 1. PDF
                if file.lower().endswith(".pdf"):
                    doc = fitz.open(file)
                    for page in doc:
                        combined_text += page.get_text() + "\n"
                
                # 2. PPTX
                elif file.lower().endswith(".pptx"):
                    prs = Presentation(file)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                combined_text += shape.text + "\n"
                
                # 3. Excel
                elif file.lower().endswith((".xlsx", ".xls")):
                    df_dict = pd.read_excel(file, sheet_name=None)
                    for sheet_name, df in df_dict.items():
                        combined_text += f"\n[시트: {sheet_name}]\n"
                        combined_text += df.to_string(index=False) + "\n"
                
                # 4. TXT
                elif file.lower().endswith(".txt"):
                    with open(file, "r", encoding="utf-8") as f:
                        combined_text += f.read() + "\n"
                
                # 5. HWP
                elif file.lower().endswith(".hwp"):
                    if olefile.isOleFile(file):
                        ole = olefile.OleFileIO(file)
                        if "PrvText" in ole.listdir():
                            combined_text += ole.openstream("PrvText").read().decode("utf-16") + "\n"
                
                # 성공적으로 읽은 파일만 리스트에 추가
                file_list.append(file)
                
            except Exception as e:
                st.sidebar.error(f"{file} 읽기 실패: {e}")
                
    return combined_text, file_list

# 데이터 로딩
knowledge_base, learned_files = load_all_documents()

# --- 4. 사이드바 (슬라이드 메뉴) 구성 ---
with st.sidebar:
    st.title("📚 학습된 문서 목록")
    st.info(f"현재 총 {len(learned_files)}개의 파일을 학습했습니다.")
    
    if learned_files:
        # 파일명을 가나다순으로 정렬하여 표시
        for i, name in enumerate(sorted(learned_files), 1):
            st.write(f"{i}. {name}")
    else:
        st.warning("학습된 문서 파일이 없습니다.")
    
    st.divider()
    if st.button("🔄 지식 새로고침"):
        st.cache_resource.clear()
        st.rerun()

# --- 5. 챗봇 UI ---
st.title("🤖 사내 통합 지식고 어시스턴트")

if "messages" not in st.session_state:
    st.session_state.messages = []

for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

if prompt := st.chat_input("문서 내용에 대해 물어보세요."):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    with st.chat_message("assistant"):
        client = Groq(api_key=GROQ_API_KEY)
        
        # Groq 모델의 컨텍스트 한도를 고려하여 텍스트 길이 조절
        context_text = knowledge_base[:40000] 
        
        messages = [
            {
                "role": "system", 
                "content": f"너는 사내 문서 전문가야. 아래 제공된 문서 내용을 바탕으로 답변해줘. 한국어로 친절하게 답변하고, 문서에 명시되지 않은 내용은 추측하지 말고 모른다고 답해줘.\n\n[문서 내용]\n{context_text}"
            }
        ]
        # 문맥 유지를 위해 최근 대화만 포함
        for m in st.session_state.messages[-5:]:
            messages.append({"role": m["role"], "content": m["content"]})

        with st.spinner("답변 생성 중..."):
            try:
                completion = client.chat.completions.create(
                    messages=messages,
                    model="llama-3.3-70b-versatile",
                    temperature=0
                )
                answer = completion.choices[0].message.content
                st.markdown(answer)
                st.session_state.messages.append({"role": "assistant", "content": answer})
            except Exception as e:
                st.error(f"오류 발생: {e}")
